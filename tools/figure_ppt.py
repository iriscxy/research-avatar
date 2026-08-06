#!/usr/bin/env python3
"""figure_ppt.py — model-figure pipeline: paper → BioRender prompt → gpt-image → editable PPT → PDF.

The mechanism (per the researcher):
  1. GENPROMPT  the fixed "expert scientific-figure designer" META_PROMPT (system) + the paper
                (user), nothing appended, via GPT chat → a BioRender-style image prompt written
                into spec.draw_prompt. The draw prompt is GENERATED from the paper, not hand-written.
  2. DRAW       that prompt drives an image model VERBATIM (nothing appended). Drawer swappable:
                --provider openai (gpt-image-1, OPENAI_API_KEY) now; gemini later. Every draw is
                ARCHIVED to iterations/<id>/round_NN.png + round_NN.prompt.txt (no overwrite).
  3. REFINE     AGENT-DRIVEN, no CLI command: the calling agent READS the drawn image, decides
                what is wrong, rewrites spec.draw_prompt itself, re-runs draw. Loop. (First-draft
                failures can't be enumerated in a fixed instruction, so there is no scripted refine.)
  4. BUILD+PDF  two paths — `buildshapes` renders a shape-spec into NATIVE editable PPT shapes
                (every element editable, crisp text); `build` lays the image as a background +
                editable label text boxes. `pdfshapes` renders the SAME shape spec to unattended
                vector PDF with headless Chrome. `pdf` converts arbitrary PPTX with a verified
                headless LibreOffice binary when one is installed.

Subcommands:
  genprompt --paper <file> --spec spec.json [--model gpt-4o]   -> writes spec.draw_prompt
  draw        spec.json  [--provider openai]                   -> <fig>.bg.png (+ iterations/)
  build       spec.json  [--img <png>]                         -> <fig>.pptx  (image bg + labels)
  buildshapes shapes.json [--out <fig>.pptx]                   -> <fig>.pptx  (fully editable shapes)
  pdfshapes   shapes.json [--out <fig>.pdf]                    -> <fig>.pdf   (unattended vector PDF)
  pdfimage    spec.json --img <png> [--out <fig>.pdf]           -> image-backed PDF
  pdf         <fig>.pptx                                        -> <fig>.pdf (LibreOffice fallback)
  all         spec.json  --paper <file>                        -> genprompt→draw→build→pdf
  emit-example                                                 -> starter spec
  (refine is agent-driven: read <fig>.bg.png, edit spec.draw_prompt, re-run draw.)
"""
import argparse
import base64
import html
import json
import os
import re
import shutil
import signal
import subprocess
import sys
import tempfile
import time
import urllib.request

# ---------- the fixed prompt-generation meta-prompt (verbatim, per the researcher) ----------
META_PROMPT = (
    "You are a professional and experienced scientific-figure designer. Carefully read the "
    "following paper content, deeply understand its core mechanism, key method, and deep-model "
    "experimental pipeline, and then generate a BioRender-style prompt for the mechanism figure."
)
def _openai_chat_raw(model, messages):
    key = os.environ.get("OPENAI_API_KEY") or sys.exit("OPENAI_API_KEY not set")
    body = json.dumps({"model": model, "messages": messages}).encode()
    req = urllib.request.Request(
        "https://api.openai.com/v1/chat/completions", data=body,
        headers={"Authorization": f"Bearer {key}", "Content-Type": "application/json"})
    with urllib.request.urlopen(req, timeout=180) as r:
        return json.load(r)["choices"][0]["message"]["content"].strip()


def _openai_chat(model, system, user):
    return _openai_chat_raw(model, [{"role": "system", "content": system},
                                    {"role": "user", "content": user}])


# NOTE: there is deliberately NO scripted "refine" step. First-draft figures fail in ways that
# cannot be enumerated in a fixed instruction, so REFINEMENT IS AGENT-DRIVEN: the calling agent
# READS the drawn image, decides what is actually wrong, rewrites `spec.draw_prompt` itself
# (Write/Edit the spec), and re-runs `draw`. See SKILL.md Stage 3.


def cmd_genprompt(args):
    paper = open(args.paper, encoding="utf-8", errors="replace").read()
    # only the meta-prompt (system) + the paper (user); nothing appended.
    prompt = _openai_chat(args.model, META_PROMPT, paper)
    if args.spec and os.path.exists(args.spec):
        spec = json.load(open(args.spec))
        spec["draw_prompt"] = prompt
        json.dump(spec, open(args.spec, "w"), ensure_ascii=False, indent=2)
        print(json.dumps({"genprompt": "written to spec.draw_prompt", "spec": args.spec,
                          "chars": len(prompt)}, ensure_ascii=False))
    print(prompt)


# ---------- 2. DRAW (swappable image model) ----------
def draw_openai(prompt, size, out_png, quality="high"):
    key = os.environ.get("OPENAI_API_KEY") or sys.exit("OPENAI_API_KEY not set")
    body = json.dumps({"model": "gpt-image-1", "prompt": prompt, "size": size,
                       "quality": quality, "n": 1}).encode()
    req = urllib.request.Request(
        "https://api.openai.com/v1/images/generations", data=body,
        headers={"Authorization": f"Bearer {key}", "Content-Type": "application/json"})
    with urllib.request.urlopen(req, timeout=300) as r:
        data = json.load(r)
    with open(out_png, "wb") as f:
        f.write(base64.b64decode(data["data"][0]["b64_json"]))
    return out_png


def draw_gemini(prompt, size, out_png, quality="high"):
    sys.exit("gemini provider not wired yet (needs GEMINI_API_KEY + google-generativeai); "
             "use --provider openai for now")


DRAWERS = {"openai": draw_openai, "gemini": draw_gemini}


def _archive_iteration(spec_path, fid, img, prompt):
    """Save every iteration (image + the exact prompt used) so nothing is overwritten."""
    d = os.path.join(os.path.dirname(os.path.abspath(spec_path)) or ".", "iterations", fid)
    os.makedirs(d, exist_ok=True)
    nums = [int(m.group(1)) for f in os.listdir(d) if (m := re.match(r"round_(\d+)\.png$", f))]
    n = max(nums, default=0) + 1
    shutil.copy(img, os.path.join(d, f"round_{n:02d}.png"))
    with open(os.path.join(d, f"round_{n:02d}.prompt.txt"), "w", encoding="utf-8") as f:
        f.write(prompt)
    return n, f"iterations/{fid}/round_{n:02d}"


def cmd_draw(args):
    spec = json.load(open(args.spec))
    if not spec.get("draw_prompt"):
        sys.exit("spec.draw_prompt is empty — run `genprompt --paper <file> --spec <spec>` first")
    out = args.out or f"{spec['figure_id']}.bg.png"
    # draw with the generated prompt VERBATIM — nothing appended.
    DRAWERS[args.provider](spec["draw_prompt"], spec.get("image_size", "1536x1024"),
                           out, spec.get("quality", "high"))
    n, saved = _archive_iteration(args.spec, spec["figure_id"], out, spec["draw_prompt"])
    print(json.dumps({"drew": out, "provider": args.provider, "iteration": n,
                      "saved": saved + ".png (+ .prompt.txt)"}))


# ---------- 3. BUILD editable PPTX ----------
def cmd_build(args):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from PIL import Image
    spec = json.load(open(args.spec))
    W, H = spec["canvas_in"]
    img = args.img or f"{spec['figure_id']}.bg.png"
    out = args.out or f"{spec['figure_id']}.pptx"
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(W), Inches(H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if os.path.exists(img):
        with Image.open(img) as image:
            image_ratio = image.width / image.height
        target_ratio = W / H
        picture = slide.shapes.add_picture(
            img, 0, 0, width=Inches(W), height=Inches(H)
        )
        # Fill the requested paper slot without stretching the GPT Image draft.
        # Prompt generation keeps critical content inside the matching safe band.
        if image_ratio > target_ratio:
            crop = (1 - target_ratio / image_ratio) / 2
            picture.crop_left = crop
            picture.crop_right = crop
        elif image_ratio < target_ratio:
            crop = (1 - image_ratio / target_ratio) / 2
            picture.crop_top = crop
            picture.crop_bottom = crop
    align = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    for lb in spec.get("labels", []):
        tb = slide.shapes.add_textbox(Inches(lb["x"] * W), Inches(lb["y"] * H),
                                      Inches(lb.get("w", 0.2) * W), Inches(0.3))
        tb.text_frame.word_wrap = True
        p = tb.text_frame.paragraphs[0]
        p.alignment = align.get(lb.get("align", "center"), PP_ALIGN.CENTER)
        run = p.add_run(); run.text = lb["text"]
        run.font.size = Pt(lb.get("size", 11)); run.font.bold = lb.get("bold", False)
        run.font.name = lb.get("font", "Times New Roman")
        if lb.get("color"):
            run.font.color.rgb = RGBColor.from_string(lb["color"])
    prs.save(out)
    print(json.dumps({"built": out, "labels": len(spec.get("labels", [])),
                      "note": "open in PowerPoint to edit label text/position by hand"}))


# ---------- 3b. BUILD as fully-editable NATIVE PPT SHAPES (all elements editable) ----------
def _hx(c):
    return c.lstrip("#")


def cmd_buildshapes(args):
    """Render a shape-spec into NATIVE PPT shapes — rounded rects / ovals / arrows / text —
    so EVERY element is editable in PowerPoint (no flat background image).
    FLAT DESIGN: solid fills only, no drop shadows on any element (shapes AND connectors set
    shadow.inherit=False), no gradients, no 3D — a clean flat schematic."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
    from pptx.oxml.ns import qn
    spec = json.load(open(args.spec))
    W, H = spec["canvas_in"]
    out = args.out or f"{spec['figure_id']}.pptx"
    prs = Presentation(); prs.slide_width, prs.slide_height = Inches(W), Inches(H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    SH = {"rounded_rect": MSO_SHAPE.ROUNDED_RECTANGLE, "rect": MSO_SHAPE.RECTANGLE,
          "oval": MSO_SHAPE.OVAL, "hexagon": MSO_SHAPE.HEXAGON, "right_arrow": MSO_SHAPE.RIGHT_ARROW}
    AL = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    IX = lambda f: Inches(f * W)
    IY = lambda f: Inches(f * H)

    def _text(tf_owner, sh):
        tf = tf_owner.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = AL.get(sh.get("align", "center"), PP_ALIGN.CENTER)
        r = p.add_run(); r.text = sh.get("text", "")
        r.font.size = Pt(sh.get("font_size", 11)); r.font.bold = sh.get("bold", False)
        r.font.name = sh.get("font", "Arial")
        if sh.get("font_color"):
            r.font.color.rgb = RGBColor.from_string(_hx(sh["font_color"]))

    def _flat(shape):
        # Flat design: NO drop shadow. shadow.inherit=False only adds an empty <a:effectLst/>, but the
        # shape's <p:style> still carries <a:effectRef idx="2"> pointing at the THEME's outer shadow,
        # and LibreOffice/soffice renders THAT on the pptx->pdf export (PowerPoint honours the override,
        # soffice does not). Neutralize the reference (idx="0") so no renderer draws a shadow.
        shape.shadow.inherit = False
        for eff in shape._element.iter(qn("a:effectRef")):
            eff.set("idx", "0")

    for sh in spec.get("shapes", []):
        k = sh["kind"]
        if k in SH:
            s = slide.shapes.add_shape(SH[k], IX(sh["x"]), IY(sh["y"]), IX(sh["w"]), IY(sh["h"]))
            if sh.get("fill"):
                s.fill.solid(); s.fill.fore_color.rgb = RGBColor.from_string(_hx(sh["fill"]))
            else:
                s.fill.background()
            if sh.get("line"):
                s.line.color.rgb = RGBColor.from_string(_hx(sh["line"])); s.line.width = Pt(sh.get("line_w", 1))
            else:
                s.line.fill.background()
            _flat(s)
            if sh.get("text"):
                _text(s, sh)
        elif k == "textbox":
            tb = slide.shapes.add_textbox(IX(sh["x"]), IY(sh["y"]), IX(sh.get("w", 0.2)), IY(sh.get("h", 0.08)))
            _text(tb, sh); _flat(tb)
        elif k in ("arrow", "line"):
            c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, IX(sh["x1"]), IY(sh["y1"]),
                                           IX(sh["x2"]), IY(sh["y2"]))
            c.line.color.rgb = RGBColor.from_string(_hx(sh.get("color", "5B6B73")))
            c.line.width = Pt(sh.get("weight", 2))
            _flat(c)  # flat — connectors must not inherit or reference the theme drop-shadow either
            if k == "arrow":
                ln = c.line._get_or_add_ln()
                ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    prs.save(out)
    print(json.dumps({"built_shapes": out, "n_shapes": len(spec.get("shapes", [])),
                      "note": "fully editable native PPT shapes — every element selectable/editable"}))


# ---------- 4. Vector PDF export ----------
def _chrome_executable():
    candidates = [
        shutil.which("google-chrome"),
        shutil.which("google-chrome-stable"),
        shutil.which("chromium"),
        shutil.which("chromium-browser"),
        "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome",
        "/Applications/Chromium.app/Contents/MacOS/Chromium",
    ]
    return next((path for path in candidates if path and os.path.isfile(path)), None)


def _svg_text(sh, x, y, w, h):
    size = float(sh.get("font_size", 11)) * 96 / 72
    color = "#" + _hx(sh.get("font_color", "222222"))
    family = html.escape(str(sh.get("font", "Arial")), quote=True)
    align = sh.get("align", "center")
    justify = {"left": "flex-start", "right": "flex-end"}.get(align, "center")
    padding = max(1.5, min(w, h) * 0.035)
    text = html.escape(str(sh.get("text", ""))).replace("\n", "<br/>")
    weight = "700" if sh.get("bold", False) else "400"
    return (
        f'<foreignObject x="{x:.3f}" y="{y:.3f}" width="{w:.3f}" height="{h:.3f}">'
        '<div xmlns="http://www.w3.org/1999/xhtml" style="box-sizing:border-box;'
        f'width:100%;height:100%;padding:{padding:.2f}px;display:flex;align-items:center;'
        f'justify-content:{justify};overflow:hidden;text-align:{align};white-space:normal;'
        f'font-family:&quot;{family}&quot;;font-size:{size:.3f}px;font-weight:{weight};'
        f'line-height:1.12;color:{color};">{text}</div></foreignObject>'
    )


def shape_spec_html(spec):
    """Render the same normalized shape spec used for PPTX as a vector SVG page."""
    W, H = map(float, spec["canvas_in"])
    width, height = W * 96, H * 96
    parts = [
        "<!doctype html><html><head><meta charset='utf-8'><style>",
        f"@page{{size:{W}in {H}in;margin:0}}",
        f"html,body{{margin:0;width:{W}in;height:{H}in;overflow:hidden;background:white}}",
        "</style></head><body>",
        f'<svg xmlns="http://www.w3.org/2000/svg" width="{W}in" height="{H}in" '
        f'viewBox="0 0 {width:.3f} {height:.3f}">',
        '<defs><marker id="arrowhead" markerWidth="7" markerHeight="7" refX="6" refY="3.5" '
        'orient="auto" markerUnits="strokeWidth"><path d="M0,0 L7,3.5 L0,7 z" fill="context-stroke"/>'
        "</marker></defs>",
        f'<rect x="0" y="0" width="{width:.3f}" height="{height:.3f}" fill="white"/>',
    ]
    for sh in spec.get("shapes", []):
        kind = sh["kind"]
        if kind in ("line", "arrow"):
            x1, y1 = float(sh["x1"]) * width, float(sh["y1"]) * height
            x2, y2 = float(sh["x2"]) * width, float(sh["y2"]) * height
            stroke = "#" + _hx(sh.get("color", "5B6B73"))
            stroke_w = float(sh.get("weight", 2)) * 96 / 72
            marker = ' marker-end="url(#arrowhead)"' if kind == "arrow" else ""
            parts.append(
                f'<line x1="{x1:.3f}" y1="{y1:.3f}" x2="{x2:.3f}" y2="{y2:.3f}" '
                f'stroke="{stroke}" stroke-width="{stroke_w:.3f}"{marker}/>'
            )
            continue
        x, y = float(sh["x"]) * width, float(sh["y"]) * height
        w, h = float(sh.get("w", .2)) * width, float(sh.get("h", .08)) * height
        if kind != "textbox":
            fill = "#" + _hx(sh["fill"]) if sh.get("fill") else "none"
            stroke = "#" + _hx(sh["line"]) if sh.get("line") else "none"
            stroke_w = float(sh.get("line_w", 1)) * 96 / 72
            attrs = f'fill="{fill}" stroke="{stroke}" stroke-width="{stroke_w:.3f}"'
            if kind in ("rounded_rect", "rect"):
                radius = min(w, h) * .12 if kind == "rounded_rect" else 0
                parts.append(
                    f'<rect x="{x:.3f}" y="{y:.3f}" width="{w:.3f}" height="{h:.3f}" '
                    f'rx="{radius:.3f}" {attrs}/>'
                )
            elif kind == "oval":
                parts.append(
                    f'<ellipse cx="{x + w/2:.3f}" cy="{y + h/2:.3f}" rx="{w/2:.3f}" '
                    f'ry="{h/2:.3f}" {attrs}/>'
                )
            elif kind == "hexagon":
                points = [(x+w*.2,y),(x+w*.8,y),(x+w,y+h/2),(x+w*.8,y+h),(x+w*.2,y+h),(x,y+h/2)]
                pts = " ".join(f"{px:.3f},{py:.3f}" for px, py in points)
                parts.append(f'<polygon points="{pts}" {attrs}/>')
            elif kind == "right_arrow":
                points = [(x,y+h*.25),(x+w*.68,y+h*.25),(x+w*.68,y),(x+w,y+h/2),
                          (x+w*.68,y+h),(x+w*.68,y+h*.75),(x,y+h*.75)]
                pts = " ".join(f"{px:.3f},{py:.3f}" for px, py in points)
                parts.append(f'<polygon points="{pts}" {attrs}/>')
        if sh.get("text"):
            parts.append(_svg_text(sh, x, y, w, h))
    parts.extend(["</svg></body></html>"])
    return "".join(parts)


def image_spec_html(spec, image_path):
    """Render the exact GPT image plus the same editable-label coordinates as the PPT."""
    W, H = map(float, spec["canvas_in"])
    encoded = base64.b64encode(open(image_path, "rb").read()).decode("ascii")
    labels = []
    for label in spec.get("labels", []):
        text = html.escape(str(label.get("text", "")))
        align = html.escape(str(label.get("align", "center")), quote=True)
        color = "#" + str(label.get("color", "222222")).lstrip("#")
        weight = "700" if label.get("bold", False) else "400"
        labels.append(
            f'<div class="label" style="left:{float(label["x"])*100:.4f}%;'
            f'top:{float(label["y"])*100:.4f}%;width:{float(label.get("w",.2))*100:.4f}%;'
            f'font-size:{float(label.get("size",11))}pt;text-align:{align};'
            f'color:{color};font-weight:{weight}">{text}</div>'
        )
    return "".join(
        [
            "<!doctype html><html><head><meta charset='utf-8'><style>",
            f"@page{{size:{W}in {H}in;margin:0}}",
            f"html,body{{margin:0;width:{W}in;height:{H}in;overflow:hidden;background:white}}",
            ".canvas{position:relative;width:100%;height:100%;overflow:hidden}",
            ".canvas img{position:absolute;inset:0;width:100%;height:100%;object-fit:cover}",
            ".label{position:absolute;height:.3in;display:flex;align-items:center;justify-content:center;",
            "box-sizing:border-box;font-family:'Times New Roman';line-height:1.05}",
            "</style></head><body><div class='canvas'>",
            f"<img src='data:image/png;base64,{encoded}'/>",
            *labels,
            "</div></body></html>",
        ]
    )


def _write_html_pdf(source_html, out, error_message):
    chrome = _chrome_executable()
    if not chrome:
        sys.exit("找不到可用的无界面 Chrome/Chromium，无法自动导出机制图 PDF。")
    out = os.path.abspath(out)
    os.makedirs(os.path.dirname(out), exist_ok=True)
    with tempfile.TemporaryDirectory(prefix="figure-pdf-") as temp_dir:
        source = os.path.join(temp_dir, "figure.html")
        # Chrome on macOS may attach com.apple.quarantine to a directly written
        # workspace PDF, after which the server itself can receive EPERM opening it.
        # Let Chrome write inside /tmp and copy only the file bytes to the deliverable.
        chrome_out = os.path.join(temp_dir, "rendered.pdf")
        with open(source, "w", encoding="utf-8") as handle:
            handle.write(source_html)
        process = subprocess.Popen(
            [chrome, "--headless=new", "--disable-gpu", "--no-pdf-header-footer",
             "--disable-background-networking", "--disable-component-update", "--disable-sync",
             "--metrics-recording-only", "--no-first-run", "--no-default-browser-check",
             "--disable-search-engine-choice-screen",
             f"--user-data-dir={os.path.join(temp_dir, 'chrome-profile')}",
             f"--print-to-pdf={chrome_out}", "file://" + source],
            stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, start_new_session=True,
        )
        deadline = time.monotonic() + 60
        stable_since = None
        previous_size = -1
        while time.monotonic() < deadline:
            if os.path.exists(chrome_out):
                current_size = os.path.getsize(chrome_out)
                if current_size >= 1000 and current_size == previous_size:
                    stable_since = stable_since or time.monotonic()
                    if time.monotonic() - stable_since >= .75:
                        break
                else:
                    stable_since = None
                    previous_size = current_size
            if process.poll() is not None:
                break
            time.sleep(.1)
        if process.poll() is None:
            try:
                os.killpg(process.pid, signal.SIGTERM)
                process.wait(timeout=3)
            except (ProcessLookupError, subprocess.TimeoutExpired):
                try:
                    os.killpg(process.pid, signal.SIGKILL)
                except ProcessLookupError:
                    pass
        if not os.path.exists(chrome_out) or os.path.getsize(chrome_out) < 1000:
            sys.exit(error_message)
        temporary_target = out + ".tmp"
        with open(chrome_out, "rb") as source_pdf, open(temporary_target, "wb") as target_pdf:
            shutil.copyfileobj(source_pdf, target_pdf)
        os.replace(temporary_target, out)
    if not os.path.exists(out) or os.path.getsize(out) < 1000:
        sys.exit("机制图 PDF 写入失败。")


def cmd_pdfshapes(args):
    spec = json.load(open(args.spec, encoding="utf-8"))
    out = os.path.abspath(args.out or f"{spec['figure_id']}.pdf")
    _write_html_pdf(shape_spec_html(spec), out, "无界面 Chrome 导出机制图 PDF 失败。")
    print(json.dumps({"pdf": out, "renderer": "headless-chrome-shape-spec", "ok": True}))


def cmd_pdfimage(args):
    spec = json.load(open(args.spec, encoding="utf-8"))
    image_path = args.img or f"{spec['figure_id']}.bg.png"
    if not os.path.exists(image_path):
        sys.exit("GPT Image 草图不存在，无法导出同视觉 PDF。")
    out = os.path.abspath(args.out or f"{spec['figure_id']}.pdf")
    _write_html_pdf(
        image_spec_html(spec, image_path),
        out,
        "无界面 Chrome 导出 GPT Image PDF 失败。",
    )
    print(json.dumps({"pdf": out, "renderer": "headless-chrome-gpt-image", "ok": True}))


def cmd_pdf(args):
    outdir = os.path.dirname(os.path.abspath(args.pptx)) or "."
    final_pdf = os.path.join(
        outdir, os.path.splitext(os.path.basename(args.pptx))[0] + ".pdf"
    )
    candidates = [shutil.which("soffice"), shutil.which("libreoffice"),
                  "/Applications/LibreOffice.app/Contents/MacOS/soffice"]
    soffice = None
    for candidate in candidates:
        if not candidate or not os.path.isfile(candidate):
            continue
        probe = subprocess.run([candidate, "--headless", "--version"], capture_output=True, timeout=10)
        if probe.returncode == 0:
            soffice = candidate
            break
    if not soffice:
        sys.exit("找不到可运行的 LibreOffice；机制图请使用 pdfshapes 进行无人值守矢量导出。")
    with tempfile.TemporaryDirectory(prefix="figure-soffice-") as temporary:
        profile = os.path.join(temporary, "profile")
        converted_dir = os.path.join(temporary, "converted")
        os.makedirs(converted_dir, exist_ok=True)
        subprocess.run(
            [soffice, "--headless", f"-env:UserInstallation=file://{profile}",
             "--convert-to", "pdf", "--outdir", converted_dir, args.pptx],
            check=True, capture_output=True,
        )
        converted_pdf = os.path.join(
            converted_dir,
            os.path.splitext(os.path.basename(args.pptx))[0] + ".pdf",
        )
        if not os.path.exists(converted_pdf) or os.path.getsize(converted_pdf) < 1000:
            sys.exit("LibreOffice 未生成有效机制图 PDF。")
        temporary_target = final_pdf + ".tmp"
        with open(converted_pdf, "rb") as source_pdf, open(temporary_target, "wb") as target_pdf:
            shutil.copyfileobj(source_pdf, target_pdf)
        os.replace(temporary_target, final_pdf)
    print(json.dumps({"pdf": final_pdf, "ok": os.path.exists(final_pdf)}))


def cmd_all(args):
    # genprompt → draw → build → pdf. The REFINE step (read image → rewrite spec.draw_prompt →
    # redraw) is agent-driven and lives BETWEEN draw and build — do it by hand, not here.
    args.out = None
    cmd_genprompt(args)
    cmd_draw(args)
    cmd_build(args)
    spec = json.load(open(args.spec))
    args.pptx = f"{spec['figure_id']}.pptx"
    cmd_pdf(args)


EXAMPLE = {
    "figure_id": "rts_fig1",
    "canvas_in": [6.5, 3.0],
    "image_size": "1536x1024",
    "quality": "high",
    "draw_prompt": "",
    "_draw_prompt_note": "leave empty — filled by `genprompt --paper <method> --spec this.json`",
    "labels": [
        {"text": "example label", "x": 0.05, "y": 0.05, "w": 0.2, "size": 10, "align": "left"}
    ]
}


def main():
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    sub = ap.add_subparsers(dest="cmd", required=True)
    gp = sub.add_parser("genprompt"); gp.add_argument("--paper", required=True)
    gp.add_argument("--spec"); gp.add_argument("--model", default="gpt-4o")
    for name in ("draw", "build", "buildshapes", "all"):
        p = sub.add_parser(name); p.add_argument("spec"); p.add_argument("--out")
        if name in ("draw", "all"):
            p.add_argument("--provider", choices=list(DRAWERS), default="openai")
        if name == "build":
            p.add_argument("--img")
        if name == "all":
            p.add_argument("--paper", required=True); p.add_argument("--model", default="gpt-4o")
    pp = sub.add_parser("pdf"); pp.add_argument("pptx")
    ps = sub.add_parser("pdfshapes"); ps.add_argument("spec"); ps.add_argument("--out")
    pi = sub.add_parser("pdfimage"); pi.add_argument("spec"); pi.add_argument("--img"); pi.add_argument("--out")
    sub.add_parser("emit-example")
    args = ap.parse_args()
    if args.cmd == "emit-example":
        print(json.dumps(EXAMPLE, ensure_ascii=False, indent=2)); return
    {"genprompt": cmd_genprompt, "draw": cmd_draw, "build": cmd_build,
     "buildshapes": cmd_buildshapes, "pdf": cmd_pdf, "pdfshapes": cmd_pdfshapes,
     "pdfimage": cmd_pdfimage,
     "all": cmd_all}[args.cmd](args)


if __name__ == "__main__":
    main()
