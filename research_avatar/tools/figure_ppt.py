#!/usr/bin/env python3
"""Build Codex-authored native PowerPoint figures and matching vector PDFs.

The authoritative input is a JSON shape specification and every rendered
paper-figure element is a native editable object.
"""
import argparse
import html
import json
import os
import shutil
import signal
import subprocess
import sys
import tempfile
import time
from pathlib import Path

# ---------- BUILD as fully-editable NATIVE PPT SHAPES ----------
def _hx(c):
    return c.lstrip("#")


SEMANTIC_SHAPE_ROLES = {"input", "operation", "output", "annotation"}
REQUIRED_SEMANTIC_SHAPE_ROLES = {"input", "operation", "output"}


def validate_native_shape_spec(spec):
    """Reject final native figures that contradict their readability contract."""
    if not isinstance(spec, dict) or not isinstance(spec.get("shapes"), list):
        raise ValueError("Native shape spec requires a shapes array.")
    if spec.get("no_text") is True:
        raise ValueError(
            "Native paper figures cannot set no_text=true; final figures require "
            "print-readable semantic labels."
        )
    try:
        contract_version = int(spec.get("semantic_contract_version") or 0)
    except (TypeError, ValueError) as exc:
        raise ValueError("semantic_contract_version must be an integer.") from exc
    if contract_version < 2:
        return spec

    text_shapes = [
        item for item in spec["shapes"]
        if isinstance(item, dict) and str(item.get("text") or "").strip()
    ]
    if len(text_shapes) < 3:
        raise ValueError(
            "Semantic figure contract requires at least three labels covering "
            "input, operation, and output."
        )
    observed_roles = set()
    observed_labels = set()
    for item in text_shapes:
        role = str(item.get("semantic_role") or "").strip()
        if role not in SEMANTIC_SHAPE_ROLES:
            raise ValueError(
                "Every text-bearing shape must declare semantic_role as input, "
                "operation, output, or annotation."
            )
        observed_roles.add(role)
        observed_labels.add(str(item["text"]).strip())
        if float(item.get("font_size", 0)) < 7:
            raise ValueError("Final paper-figure labels must be at least 7 pt.")
    required_roles = {
        str(item).strip()
        for item in spec.get(
            "required_semantic_roles", sorted(REQUIRED_SEMANTIC_SHAPE_ROLES)
        )
        if str(item).strip()
    }
    missing_roles = required_roles - observed_roles
    if missing_roles:
        raise ValueError(
            "Semantic figure contract is missing roles: "
            + ", ".join(sorted(missing_roles))
        )
    required_labels = {
        str(item).strip()
        for item in spec.get("required_labels", [])
        if str(item).strip()
    }
    missing_labels = required_labels - observed_labels
    if missing_labels:
        raise ValueError(
            "Native figure is missing required labels: "
            + ", ".join(sorted(missing_labels))
        )
    return spec


def cmd_buildshapes(args):
    """Render a shape-spec into NATIVE PPT shapes — rounded rects / ovals / arrows / text —
    so EVERY element is editable in PowerPoint (no flat background image).
    FLAT DESIGN: solid fills only, no drop shadows on any element (shapes AND connectors set
    shadow.inherit=False), no gradients, no 3D — a clean flat schematic."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
    from pptx.oxml.ns import qn
    spec = validate_native_shape_spec(
        json.loads(Path(args.spec).read_text(encoding="utf-8"))
    )
    W, H = spec["canvas_in"]
    out = args.out or f"{spec['figure_id']}.pptx"
    prs = Presentation(); prs.slide_width, prs.slide_height = Inches(W), Inches(H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    SH = {"rounded_rect": MSO_SHAPE.ROUNDED_RECTANGLE, "rect": MSO_SHAPE.RECTANGLE,
          "oval": MSO_SHAPE.OVAL, "hexagon": MSO_SHAPE.HEXAGON, "right_arrow": MSO_SHAPE.RIGHT_ARROW}
    AL = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    IX = lambda f: Inches(f * W)
    IY = lambda f: Inches(f * H)

    def _text(tf_owner, sh):
        tf = tf_owner.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = AL.get(sh.get("align", "center"), PP_ALIGN.CENTER)
        r = p.add_run(); r.text = sh.get("text", "")
        r.font.size = Pt(sh.get("font_size", 11)); r.font.bold = sh.get("bold", False)
        r.font.name = sh.get("font", "Arial")
        if sh.get("font_color"):
            r.font.color.rgb = RGBColor.from_string(_hx(sh["font_color"]))

    def _flat(shape):
        # Flat design: NO drop shadow. shadow.inherit=False only adds an empty <a:effectLst/>, but the
        # shape's <p:style> still carries <a:effectRef idx="2"> pointing at the THEME's outer shadow,
        # and LibreOffice/soffice renders THAT on the pptx->pdf export (PowerPoint honours the override,
        # soffice does not). Neutralize the reference (idx="0") so no renderer draws a shadow.
        shape.shadow.inherit = False
        for eff in shape._element.iter(qn("a:effectRef")):
            eff.set("idx", "0")

    for sh in spec.get("shapes", []):
        k = sh["kind"]
        if k in SH:
            s = slide.shapes.add_shape(SH[k], IX(sh["x"]), IY(sh["y"]), IX(sh["w"]), IY(sh["h"]))
            if sh.get("fill"):
                s.fill.solid(); s.fill.fore_color.rgb = RGBColor.from_string(_hx(sh["fill"]))
            else:
                s.fill.background()
            if sh.get("line"):
                s.line.color.rgb = RGBColor.from_string(_hx(sh["line"])); s.line.width = Pt(sh.get("line_w", 1))
            else:
                s.line.fill.background()
            _flat(s)
            if sh.get("text"):
                _text(s, sh)
        elif k == "textbox":
            tb = slide.shapes.add_textbox(IX(sh["x"]), IY(sh["y"]), IX(sh.get("w", 0.2)), IY(sh.get("h", 0.08)))
            _text(tb, sh); _flat(tb)
        elif k in ("arrow", "line"):
            c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, IX(sh["x1"]), IY(sh["y1"]),
                                           IX(sh["x2"]), IY(sh["y2"]))
            c.line.color.rgb = RGBColor.from_string(_hx(sh.get("color", "5B6B73")))
            c.line.width = Pt(sh.get("weight", 2))
            if sh.get("dash"):
                ln = c.line._get_or_add_ln()
                for existing in list(ln.findall(qn("a:prstDash"))):
                    ln.remove(existing)
                ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
            _flat(c)  # flat — connectors must not inherit or reference the theme drop-shadow either
            if k == "arrow":
                ln = c.line._get_or_add_ln()
                ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    prs.save(out)
    print(json.dumps({"built_shapes": out, "n_shapes": len(spec.get("shapes", [])),
                      "note": "fully editable native PPT shapes — every element selectable/editable"}))


# ---------- 4. Vector PDF export ----------
def _chrome_executable():
    candidates = [
        shutil.which("google-chrome"),
        shutil.which("google-chrome-stable"),
        shutil.which("chromium"),
        shutil.which("chromium-browser"),
        "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome",
        "/Applications/Chromium.app/Contents/MacOS/Chromium",
    ]
    return next((path for path in candidates if path and os.path.isfile(path)), None)


def _svg_text(sh, x, y, w, h):
    size = float(sh.get("font_size", 11)) * 96 / 72
    color = "#" + _hx(sh.get("font_color", "222222"))
    family = html.escape(str(sh.get("font", "Arial")), quote=True)
    align = sh.get("align", "center")
    justify = {"left": "flex-start", "right": "flex-end"}.get(align, "center")
    padding = max(1.5, min(w, h) * 0.035)
    text = html.escape(str(sh.get("text", ""))).replace("\n", "<br/>")
    weight = "700" if sh.get("bold", False) else "400"
    return (
        f'<foreignObject x="{x:.3f}" y="{y:.3f}" width="{w:.3f}" height="{h:.3f}">'
        '<div xmlns="http://www.w3.org/1999/xhtml" style="box-sizing:border-box;'
        f'width:100%;height:100%;padding:{padding:.2f}px;display:flex;align-items:center;'
        f'justify-content:{justify};overflow:hidden;text-align:{align};white-space:normal;'
        f'font-family:&quot;{family}&quot;;font-size:{size:.3f}px;font-weight:{weight};'
        f'line-height:1.12;color:{color};">{text}</div></foreignObject>'
    )


def shape_spec_html(spec):
    """Render the same normalized shape spec used for PPTX as a vector SVG page."""
    W, H = map(float, spec["canvas_in"])
    width, height = W * 96, H * 96
    parts = [
        "<!doctype html><html><head><meta charset='utf-8'><style>",
        f"@page{{size:{W}in {H}in;margin:0}}",
        f"html,body{{margin:0;width:{W}in;height:{H}in;overflow:hidden;background:white}}",
        "</style></head><body>",
        f'<svg xmlns="http://www.w3.org/2000/svg" width="{W}in" height="{H}in" '
        f'viewBox="0 0 {width:.3f} {height:.3f}">',
        '<defs><marker id="arrowhead" markerWidth="7" markerHeight="7" refX="6" refY="3.5" '
        'orient="auto" markerUnits="strokeWidth"><path d="M0,0 L7,3.5 L0,7 z" fill="context-stroke"/>'
        "</marker></defs>",
        f'<rect x="0" y="0" width="{width:.3f}" height="{height:.3f}" fill="white"/>',
    ]
    for sh in spec.get("shapes", []):
        kind = sh["kind"]
        if kind in ("line", "arrow"):
            x1, y1 = float(sh["x1"]) * width, float(sh["y1"]) * height
            x2, y2 = float(sh["x2"]) * width, float(sh["y2"]) * height
            stroke = "#" + _hx(sh.get("color", "5B6B73"))
            stroke_w = float(sh.get("weight", 2)) * 96 / 72
            marker = ' marker-end="url(#arrowhead)"' if kind == "arrow" else ""
            dash = f' stroke-dasharray="{stroke_w * 4:.3f},{stroke_w * 3:.3f}"' if sh.get("dash") else ""
            parts.append(
                f'<line x1="{x1:.3f}" y1="{y1:.3f}" x2="{x2:.3f}" y2="{y2:.3f}" '
                f'stroke="{stroke}" stroke-width="{stroke_w:.3f}"{dash}{marker}/>'
            )
            continue
        x, y = float(sh["x"]) * width, float(sh["y"]) * height
        w, h = float(sh.get("w", .2)) * width, float(sh.get("h", .08)) * height
        if kind != "textbox":
            fill = "#" + _hx(sh["fill"]) if sh.get("fill") else "none"
            stroke = "#" + _hx(sh["line"]) if sh.get("line") else "none"
            stroke_w = float(sh.get("line_w", 1)) * 96 / 72
            attrs = f'fill="{fill}" stroke="{stroke}" stroke-width="{stroke_w:.3f}"'
            if kind in ("rounded_rect", "rect"):
                radius = min(w, h) * .12 if kind == "rounded_rect" else 0
                parts.append(
                    f'<rect x="{x:.3f}" y="{y:.3f}" width="{w:.3f}" height="{h:.3f}" '
                    f'rx="{radius:.3f}" {attrs}/>'
                )
            elif kind == "oval":
                parts.append(
                    f'<ellipse cx="{x + w/2:.3f}" cy="{y + h/2:.3f}" rx="{w/2:.3f}" '
                    f'ry="{h/2:.3f}" {attrs}/>'
                )
            elif kind == "hexagon":
                points = [(x+w*.2,y),(x+w*.8,y),(x+w,y+h/2),(x+w*.8,y+h),(x+w*.2,y+h),(x,y+h/2)]
                pts = " ".join(f"{px:.3f},{py:.3f}" for px, py in points)
                parts.append(f'<polygon points="{pts}" {attrs}/>')
            elif kind == "right_arrow":
                points = [(x,y+h*.25),(x+w*.68,y+h*.25),(x+w*.68,y),(x+w,y+h/2),
                          (x+w*.68,y+h),(x+w*.68,y+h*.75),(x,y+h*.75)]
                pts = " ".join(f"{px:.3f},{py:.3f}" for px, py in points)
                parts.append(f'<polygon points="{pts}" {attrs}/>')
        if sh.get("text"):
            parts.append(_svg_text(sh, x, y, w, h))
    parts.extend(["</svg></body></html>"])
    return "".join(parts)


def _write_html_pdf(source_html, out, error_message):
    chrome = _chrome_executable()
    if not chrome:
        sys.exit("Cannot find a usable headless Chrome or Chromium, cannot automatically export mechanism diagram PDF.")
    out = os.path.abspath(out)
    os.makedirs(os.path.dirname(out), exist_ok=True)
    with tempfile.TemporaryDirectory(prefix="figure-pdf-") as temp_dir:
        source = os.path.join(temp_dir, "figure.html")
        # Chrome on macOS may attach com.apple.quarantine to a directly written
        # workspace PDF, after which the server itself can receive EPERM opening it.
        # Let Chrome write inside /tmp and copy only the file bytes to the deliverable.
        chrome_out = os.path.join(temp_dir, "rendered.pdf")
        with open(source, "w", encoding="utf-8") as handle:
            handle.write(source_html)
        process = subprocess.Popen(
            [chrome, "--headless=new", "--disable-gpu", "--no-pdf-header-footer",
             "--disable-background-networking", "--disable-component-update", "--disable-sync",
             "--metrics-recording-only", "--no-first-run", "--no-default-browser-check",
             "--disable-search-engine-choice-screen",
             f"--user-data-dir={os.path.join(temp_dir, 'chrome-profile')}",
             f"--print-to-pdf={chrome_out}", "file://" + source],
            stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, start_new_session=True,
        )
        deadline = time.monotonic() + 60
        stable_since = None
        previous_size = -1
        while time.monotonic() < deadline:
            if os.path.exists(chrome_out):
                current_size = os.path.getsize(chrome_out)
                if current_size >= 1000 and current_size == previous_size:
                    stable_since = stable_since or time.monotonic()
                    if time.monotonic() - stable_since >= .75:
                        break
                else:
                    stable_since = None
                    previous_size = current_size
            if process.poll() is not None:
                break
            time.sleep(.1)
        if process.poll() is None:
            try:
                os.killpg(process.pid, signal.SIGTERM)
                process.wait(timeout=3)
            except (ProcessLookupError, subprocess.TimeoutExpired):
                try:
                    os.killpg(process.pid, signal.SIGKILL)
                except ProcessLookupError:
                    pass
        if not os.path.exists(chrome_out) or os.path.getsize(chrome_out) < 1000:
            sys.exit(error_message)
        temporary_target = out + ".tmp"
        with open(chrome_out, "rb") as source_pdf, open(temporary_target, "wb") as target_pdf:
            shutil.copyfileobj(source_pdf, target_pdf)
        os.replace(temporary_target, out)
    if not os.path.exists(out) or os.path.getsize(out) < 1000:
        sys.exit("Failed to write mechanism diagram PDF.")


def cmd_pdfshapes(args):
    spec = validate_native_shape_spec(
        json.loads(Path(args.spec).read_text(encoding="utf-8"))
    )
    out = os.path.abspath(args.out or f"{spec['figure_id']}.pdf")
    _write_html_pdf(shape_spec_html(spec), out, "Chrome export of the mechanism diagram PDF failed without UI.")
    print(json.dumps({"pdf": out, "renderer": "headless-chrome-shape-spec", "ok": True}))


def cmd_pdf(args):
    outdir = os.path.dirname(os.path.abspath(args.pptx)) or "."
    final_pdf = os.path.join(
        outdir, os.path.splitext(os.path.basename(args.pptx))[0] + ".pdf"
    )
    candidates = [shutil.which("soffice"), shutil.which("libreoffice"),
                  "/Applications/LibreOffice.app/Contents/MacOS/soffice"]
    soffice = None
    for candidate in candidates:
        if not candidate or not os.path.isfile(candidate):
            continue
        probe = subprocess.run([candidate, "--headless", "--version"], capture_output=True, timeout=10)
        if probe.returncode == 0:
            soffice = candidate
            break
    if not soffice:
        sys.exit("Cannot find a runnable LibreOffice; mechanism diagrams should be exported as unattended vectors using pdfshapes.")
    with tempfile.TemporaryDirectory(prefix="figure-soffice-") as temporary:
        profile = os.path.join(temporary, "profile")
        converted_dir = os.path.join(temporary, "converted")
        os.makedirs(converted_dir, exist_ok=True)
        subprocess.run(
            [soffice, "--headless", f"-env:UserInstallation=file://{profile}",
             "--convert-to", "pdf", "--outdir", converted_dir, args.pptx],
            check=True, capture_output=True,
        )
        converted_pdf = os.path.join(
            converted_dir,
            os.path.splitext(os.path.basename(args.pptx))[0] + ".pdf",
        )
        if not os.path.exists(converted_pdf) or os.path.getsize(converted_pdf) < 1000:
            sys.exit("LibreOffice No valid mechanism diagram PDF generated.")
        temporary_target = final_pdf + ".tmp"
        with open(converted_pdf, "rb") as source_pdf, open(temporary_target, "wb") as target_pdf:
            shutil.copyfileobj(source_pdf, target_pdf)
        os.replace(temporary_target, final_pdf)
    print(json.dumps({"pdf": final_pdf, "ok": os.path.exists(final_pdf)}))


EXAMPLE = {
    "figure_id": "rts_fig1",
    "canvas_in": [6.5, 3.0],
    "shapes": [
        {"kind": "rounded_rect", "x": 0.08, "y": 0.25, "w": 0.22, "h": 0.3,
         "fill": "EAF4F8", "line": "203746", "line_w": 1,
         "text": "Input", "font_size": 9, "bold": True,
         "font_color": "203746", "align": "center"},
        {"kind": "arrow", "x1": 0.32, "y1": 0.4, "x2": 0.52, "y2": 0.4,
         "color": "203746", "weight": 1.5},
        {"kind": "hexagon", "x": 0.55, "y": 0.22, "w": 0.3, "h": 0.36,
         "fill": "EAF7F5", "line": "087F74", "line_w": 1,
         "text": "Operation", "font_size": 9, "bold": True,
         "font_color": "087F74", "align": "center"}
    ]
}


def main():
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    sub = ap.add_subparsers(dest="cmd", required=True)
    bs = sub.add_parser("buildshapes"); bs.add_argument("spec"); bs.add_argument("--out")
    pp = sub.add_parser("pdf"); pp.add_argument("pptx")
    ps = sub.add_parser("pdfshapes"); ps.add_argument("spec"); ps.add_argument("--out")
    sub.add_parser("emit-example")
    args = ap.parse_args()
    if args.cmd == "emit-example":
        print(json.dumps(EXAMPLE, ensure_ascii=False, indent=2)); return
    {"buildshapes": cmd_buildshapes, "pdf": cmd_pdf,
     "pdfshapes": cmd_pdfshapes}[args.cmd](args)


if __name__ == "__main__":
    main()
